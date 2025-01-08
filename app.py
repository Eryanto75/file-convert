import streamlit as st
from werkzeug.utils import secure_filename
import os
from fpdf import FPDF
from docx import Document
import pandas as pd
from pptx import Presentation

UPLOAD_FOLDER = 'uploads/'

# Membuat folder upload jika belum ada
if not os.path.exists(UPLOAD_FOLDER):
    os.makedirs(UPLOAD_FOLDER)

def convert_docx_to_pdf(file, filename):
    doc = Document(file)
    pdf = FPDF()
    pdf.add_page()
    pdf.set_auto_page_break(auto=True, margin=15)
    pdf.set_font("Arial", size=12)
    for para in doc.paragraphs:
        pdf.multi_cell(0, 10, para.text)
    output_path = os.path.join(UPLOAD_FOLDER, filename.replace(".docx", ".pdf"))
    pdf.output(output_path)
    return output_path

def convert_xlsx_to_pdf(file, filename):
    df = pd.read_excel(file)
    pdf = FPDF()
    pdf.add_page()
    pdf.set_auto_page_break(auto=True, margin=15)
    pdf.set_font("Arial", size=12)
    for index, row in df.iterrows():
        row_str = ' '.join([str(item) for item in row])
        pdf.multi_cell(0, 10, row_str)
    output_path = os.path.join(UPLOAD_FOLDER, filename.replace(".xlsx", ".pdf"))
    pdf.output(output_path)
    return output_path

def convert_pptx_to_pdf(file, filename):
    presentation = Presentation(file)
    pdf = FPDF()
    pdf.set_auto_page_break(auto=True, margin=15)
    pdf.set_font("Arial", size=12)
    for slide in presentation.slides:
        pdf.add_page()
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                pdf.multi_cell(0, 10, shape.text)
    output_path = os.path.join(UPLOAD_FOLDER, filename.replace(".pptx", ".pdf"))
    pdf.output(output_path)
    return output_path

st.title("Convert Document to PDF format")

uploaded_file = st.file_uploader("Unggah file Anda", type=["docx", "xlsx", "pptx"])

if uploaded_file is not None:
    filename = secure_filename(uploaded_file.name)
    if filename.endswith(".docx"):
        file_path = convert_docx_to_pdf(uploaded_file, filename)
    elif filename.endswith(".xlsx"):
        file_path = convert_xlsx_to_pdf(uploaded_file, filename)
    elif filename.endswith(".pptx"):
        file_path = convert_pptx_to_pdf(uploaded_file, filename)
    st.success(f"File {filename} berhasil diunggah dan dikonversi.")
    st.download_button(label="Unduh PDF", data=open(file_path, 'rb'), file_name=filename.replace(filename.split('.')[-1], 'pdf'))